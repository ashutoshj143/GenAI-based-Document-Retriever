import os
import streamlit as st
from docx import Document
import PyPDF2
import pptx
import openpyxl
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
import google.generativeai as genai

# Extraction functions 

def extract_text_from_pdf(file_path):
    text = ''
    with open(file_path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        for page_num in range(len(reader.pages)):
            text += reader.pages[page_num].extract_text()
    return text

def extract_text_from_pptx(file_path):
    text = ''
    prs = pptx.Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + '\n'
    return text

def extract_text_from_xlsx(file_path):
    text = ''
    wb = openpyxl.load_workbook(file_path)
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows():
            for cell in row:
                if cell.value:
                    text += str(cell.value) + ' '
            text += '\n'
    return text

def extract_text_from_docx(file_path):
    text = ''
    doc = Document(file_path)
    for para in doc.paragraphs:
        text += para.text + '\n'
    return text

def extract_text_from_file(file_path):
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext == '.pptx':
        return extract_text_from_pptx(file_path)
    elif ext == '.xlsx':
        return extract_text_from_xlsx(file_path)
    elif ext == '.docx':
        return extract_text_from_docx(file_path)
    else:
        return ""

# Chunking 

def chunk_text(text):
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=50,
        length_function=len,
    )
    chunks = text_splitter.split_text(text)
    return chunks

# Embeddings

def create_and_store_embeddings(chunks, metadatas):
    try:
        embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
        vectorstore = FAISS.from_texts(chunks, embeddings, metadatas=metadatas)
        return vectorstore
    except Exception as e:
        st.error(f"Error creating embeddings: {str(e)}")
        return None

# Model initialization

def initialize_model():
    genai.configure(api_key="AIzaSyDuj2Nvl9jf7Zz7z6yiGymmmWlfb9zo6lM")
    model = genai.GenerativeModel('gemini-pro')
    return model

# Answer query function
def answer_query(query, context, model):
    prompt = f"Context: {context}\n\nQuestion: {query}\n\nAnswer:"
    response = model.generate_content(prompt)
    response_text = response.text.lower()
    
    out_of_context_phrases = [
        "out of context",
        "does not mention",
        "doesn't mention",
        "not mentioned",
        "no information",
        "cannot answer",
        "unable to answer",
        "no mention of",
        "not provided in the context",
        "not enough information"
    ]
    
    if any(phrase in response_text for phrase in out_of_context_phrases):
        raise ValueError("Out of context or insufficient information")
    
    return response.text

# Streamlit UI enhancements
st.set_page_config(page_title="Tata Elxsi AI Document Assistant", page_icon="📚", layout="wide")

# Custom CSS
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Poppins:wght@300;400;600&display=swap');
    
    body {
        font-family: 'Poppins', sans-serif;
        background-color: #f0f8ff;
        color: #333;
    }
    .main {
        background: linear-gradient(135deg, #e6f3ff, #f0f8ff);
    }
    .stApp {
        max-width: 1200px;
        margin: 0 auto;
    }
    .st-bw {
        background-color: #ffffff;
        border-radius: 15px;
        padding: 30px;
        box-shadow: 0 10px 20px rgba(0, 119, 181, 0.1);
    }
    .st-er {
        font-weight: bold;
        color: #ff4b4b;
    }
    .stTextInput > div > div > input {
        background-color: #f0f8ff;
        border: 2px solid #0077B5;
        border-radius: 10px;
        padding: 10px;
        font-size: 16px;
    }
    .stButton>button {
        background: linear-gradient(45deg, #0077B5, #00a0dc);
        color: white;
        border-radius: 25px;
        padding: 12px 30px;
        font-weight: 600;
        font-size: 16px;
        transition: all 0.3s ease;
        border: none;
        box-shadow: 0 5px 15px rgba(0, 119, 181, 0.3);
    }
    .stButton>button:hover {
        background: linear-gradient(45deg, #005582, #0077B5);
        transform: translateY(-3px);
        box-shadow: 0 8px 20px rgba(0, 119, 181, 0.4);
    }
    .css-1v0mbdj.etr89bj1 {
        display: block;
        margin-left: auto;
        margin-right: auto;
        min-width: 180px;
    }
    .header {
        display: flex;
        align-items: center;
        background: linear-gradient(90deg, #003366, #0077B5);
        padding: 20px;
        border-radius: 15px;
        margin-bottom: 30px;
        box-shadow: 0 10px 30px rgba(0, 51, 102, 0.2);
    }
    .header img {
        width: 120px;
        margin-right: 20px;
    }
    .header h1 {
        color: white;
        margin: 0;
        font-size: 32px;
        font-weight: 600;
        text-shadow: 2px 2px 4px rgba(0,0,0,0.3);
    }
    .stats-container {
        display: flex;
        justify-content: space-around;
        margin-bottom: 30px;
    }
    .stat-box {
        background: linear-gradient(135deg, #ffffff, #f0f8ff);
        border-radius: 15px;
        padding: 20px;
        text-align: center;
        box-shadow: 0 8px 16px rgba(0, 119, 181, 0.1);
        transition: all 0.3s ease;
        width: 30%;
    }
    .stat-box:hover {
        transform: translateY(-5px);
        box-shadow: 0 12px 24px rgba(0, 119, 181, 0.2);
    }
    .stat-box h3 {
        margin: 0;
        color: #003366;
        font-size: 18px;
        font-weight: 600;
    }
    .stat-box p {
        font-size: 28px;
        font-weight: bold;
        margin: 10px 0 0 0;
        color: #0077B5;
    }
    .stExpander {
        border: none;
        box-shadow: 0 4px 6px rgba(0, 119, 181, 0.1);
        border-radius: 10px;
        margin-bottom: 10px;
    }
    .stExpander > div > div > div > div {
        background-color: #f0f8ff;
        border-radius: 10px 10px 0 0;
    }
    .stExpander > div > div > div > div:hover {
        background-color: #e6f3ff;
    }
    .stAlert {
        border-radius: 10px;
        border: none;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
    }
</style>
""", unsafe_allow_html=True)

def get_file_content_as_bytes(file_path):
    """Reads the file content and returns it as bytes."""
    with open(file_path, "rb") as f:
        file_content = f.read()
    return file_content

@st.cache_data
def load_documents(root_dir):
    all_documents = []
    
    for root, _, files in os.walk(root_dir):
        for file in files:
            file_path = os.path.join(root, file)
            text = extract_text_from_file(file_path)
            if text.strip():  # Ensure there's actual text in the file
                all_documents.append({
                    'content': text,
                    'metadata': {'source': file_path}
                })
    
    return all_documents

@st.cache_resource
def load_or_create_vectorstore(root_dir):
    vectorstore_path = "faiss_vectorstore"
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    
    if os.path.exists(vectorstore_path):
        try:
            vectorstore = FAISS.load_local(vectorstore_path, embeddings, allow_dangerous_deserialization=True)
            st.success("Loaded existing vector store.")
        except Exception as e:
            st.warning(f"Failed to load existing vector store: {str(e)}. Creating a new one.")
            vectorstore = None
    else:
        vectorstore = None
    
    if vectorstore is None:
        all_documents = load_documents(root_dir)
        chunks = []
        metadatas = []
        for doc in all_documents:
            doc_chunks = chunk_text(doc['content'])
            chunks.extend(doc_chunks)
            metadatas.extend([doc['metadata']] * len(doc_chunks))
        
        vectorstore = create_and_store_embeddings(chunks, metadatas)
        if vectorstore:
            vectorstore.save_local(vectorstore_path)
            st.success("Created and saved new vector store.")
    
    return vectorstore

# ... (keep all the previous code)

def process_query():
    if not st.session_state.query:
        st.warning("Please enter a question.")
        return

    with st.spinner('🧠 Thinking...'):
        root_dir = "C:\\Users\\41827\\Downloads\\MCB\\Broadcast & Media"  # Update this path
        vectorstore = load_or_create_vectorstore(root_dir)
        if vectorstore is None:
            st.error("Failed to load or create vector store. Please check your data.")
            return

        model = initialize_model()
        retriever = vectorstore.as_retriever(search_kwargs={"k": 10})
        st.session_state.relevant_docs = retriever.get_relevant_documents(st.session_state.query)
        
        try:
            context = "\n".join([doc.page_content for doc in st.session_state.relevant_docs])
            st.session_state.result = answer_query(st.session_state.query, context, model)
        except ValueError as e:
            st.session_state.result = str(e)
            st.session_state.relevant_docs = []  # Clear relevant docs

    # Clear the input after processing
    st.session_state.query = ''
def main():
    # Header with Tata Elxsi logo
    st.markdown("""
    <div class="header">
        <img src="https://medicalalley.org/wp-content/uploads/2020/10/tata-elxsi-300x243.jpg" alt="Tata Elxsi Logo">
        <h1>FileWise - AI Document Assistant</h1>
    </div>
    """, unsafe_allow_html=True)

    # Document Stats
    st.markdown("""
    <div class="stats-container">
        <div class="stat-box">
            <h3>Documents Indexed</h3>
            <p>1,234</p>
        </div>
        <div class="stat-box">
            <h3>Total Pages</h3>
            <p>5,678</p>
        </div>
        <div class="stat-box">
            <h3>File Types</h3>
            <p>4</p>
        </div>
    </div>
    """, unsafe_allow_html=True)

    # Initialize session state
    if 'query' not in st.session_state:
        st.session_state.query = ''
    if 'result' not in st.session_state:
        st.session_state.result = ''
    if 'relevant_docs' not in st.session_state:
        st.session_state.relevant_docs = []

    st.markdown("### 🔍 Ask me anything about your documents!")
    query = st.text_input("Enter your question:", value="", key="query_input")

    # Check for Enter key press
    if query != st.session_state.query:
        st.session_state.query = query
        process_query()

    if st.button("🚀 Get Answer"):
        st.session_state.query = query
        process_query()

    if st.session_state.result:
        st.markdown("### 🎉 Here's what I found:")
        st.info(st.session_state.result)

        # Only display documents if there are relevant docs (which are cleared if out of context)
        if st.session_state.relevant_docs:
            st.markdown("### 📑 Source Documents")
            for i, doc in enumerate(st.session_state.relevant_docs):
                with st.expander(f"Document {i+1}"):
                    st.write(f"**Content:** {doc.page_content}")
                    st.write(f"**Source:** {doc.metadata['source']}")
                    
                    # Add download button for the entire file
                    download_button = st.download_button(
                        label=f"📥 Download Document {i+1}", 
                        data=get_file_content_as_bytes(doc.metadata['source']), 
                        file_name=os.path.basename(doc.metadata['source']),
                        key=f"download_{i}"
                    )
                    
                    if download_button:
                        st.success(f"✅ File '{os.path.basename(doc.metadata['source'])}' successfully downloaded!")
        else:
            st.warning("No documents available as this query is out of context.")

    st.markdown("---")
    st.markdown("### 💡 Pro Tips")
    st.info("""
    - Be specific in your questions for more accurate answers.
    - You can ask about multiple documents at once.
    - Try asking for summaries or comparisons between documents.
    """)

if __name__ == "__main__":
    main()